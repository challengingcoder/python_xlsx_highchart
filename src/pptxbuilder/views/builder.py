import os
import tempfile
from base64 import b64decode
from gettext import gettext as _
from io import BytesIO

import numpy as np
import pandas as pd
from flask import Blueprint, request, abort, session, redirect, url_for, render_template
from flask import jsonify
from flask import send_file
from flask.views import MethodView
from pptx import Presentation

from pptxbuilder.constants import *
from pptxbuilder.excel_parser import parse, UnsupportedFileException, IncompatibleExcelException
from pptxbuilder.helper import str_chart_type, str_series_opt
from pptxbuilder.pptx_helper import (
    add_bar_chart,
    add_column_chart,
    add_line_chart,
    add_doughnut_chart,
    PPT_CUSTOM_LABEL_SEPERATOR)
from pptxbuilder.util import random_alphanum

builder_bp = Blueprint('builder', __name__)

builder_bp.add_app_template_global(SERIES_BY_CATEGORIES, name='SERIES_BY_CATEGORIES')
builder_bp.add_app_template_global(SERIES_BY_OPTIONS, name='SERIES_BY_OPTIONS')
builder_bp.add_app_template_global(str_series_opt)

builder_bp.add_app_template_global(CHART_TYPE_BAR, name='CHART_TYPE_BAR')
builder_bp.add_app_template_global(CHART_TYPE_COLUMN, name='CHART_TYPE_COLUMN')
builder_bp.add_app_template_global(CHART_TYPE_LINE, name='CHART_TYPE_LINE')
builder_bp.add_app_template_global(CHART_TYPE_PIE, name='CHART_TYPE_PIE')
builder_bp.add_app_template_global(str_chart_type)

this_dir = os.path.abspath(os.path.dirname(os.path.realpath(__file__)))


class Create(MethodView):
    def post(self):
        xls_str = request.form.get('file')

        try:
            xls_str = xls_str.split('base64,')[1]
            xls_data = b64decode(xls_str)
        except IndexError:
            abort(406, 'Invalid form data')

        save_path = os.path.join(tempfile.gettempdir(), random_alphanum(10))
        with open(save_path, 'wb') as f:
            f.write(xls_data)
            f.close()

        try:
            bundle = parse(save_path)
        except UnsupportedFileException:
            abort(406, 'Unsupported file format')
        except IncompatibleExcelException as ex:
            abort(406, str(ex))
        finally:
            os.remove(save_path)

        tables = bundle.to_dict()['tables'][:]
        for table in tables:
            table['view_settings'] = {
                'cross_break': 0,
                'chart_type': CHART_TYPE_BAR,
                'series_by': SERIES_BY_CATEGORIES
            }

        session['tables'] = tables
        session['builder_flag'] = 1

        return redirect(url_for('builder.index'))


class Index(MethodView):
    def get(self):
        if 'builder_flag' not in session:
            return redirect(url_for('home.index'))

        return redirect(url_for('builder.table', table_index=0))


class TableAtIndex(MethodView):
    def get(self, table_index):
        if 'builder_flag' not in session:
            return redirect(url_for('home.index'))

        tables = session['tables']

        try:
            table = session['tables'][table_index]
        except IndexError:
            abort(404)

        if request.args.get('json'):
            return jsonify(table)

        return render_template('builder/table.jinja2',
                               tables=tables,
                               table=table,
                               index=table_index)

    def post(self, table_index):
        cross_break = request.form.get('cross-break', type=int)
        chart_type = request.form.get('chart-type', type=int)
        series_by = request.form.get('series-by', type=int)

        if chart_type not in [CHART_TYPE_BAR, CHART_TYPE_COLUMN, CHART_TYPE_LINE, CHART_TYPE_PIE]:
            abort(400, _('Invalid chart type'))

        if series_by not in [SERIES_BY_CATEGORIES, SERIES_BY_OPTIONS]:
            abort(400, _('Invalid series option'))

        try:
            session['tables'][table_index]['sections'][cross_break]
        except IndexError:
            abort(400, _('Invalid crossbreak'))

        try:
            table = session['tables'][table_index]
        except IndexError:
            abort(404)

        table['view_settings']['cross_break'] = cross_break
        table['view_settings']['chart_type'] = chart_type
        table['view_settings']['series_by'] = series_by

        return redirect(url_for('builder.table', table_index=table_index))


class Export(MethodView):
    def get(self):
        if 'builder_flag' not in session:
            return redirect(url_for('home.index'))

        template_path = os.path.join(this_dir, '..', 'assets', 'template.pptx')
        prs = Presentation(template_path)
        tables = session['tables']

        # add slides table by table
        for table in tables:

            slide = prs.slides.add_slide(prs.slide_layouts[0])

            # table information
            question = table['question']
            options = table['options']
            cross_break = table['view_settings']['cross_break']
            section = table['sections'][cross_break]
            chart_type = table['view_settings']['chart_type']
            series_by = table['view_settings']['series_by']
            sub_title_text = _('Crossbreak: {}'.format(section['name']))
            reverse_colors = (series_by == SERIES_BY_OPTIONS)

            # Build pre chart data to render
            pd_list = {}
            for c in section['categories']:
                pd_list[c['name']] = c['data']

            data_frame = pd.DataFrame(pd_list, options)

            if series_by == SERIES_BY_OPTIONS:
                data_frame = data_frame.transpose()

            # fix table
            data_frame = (data_frame / 100).replace(np.NaN, '')

            if chart_type == CHART_TYPE_BAR:
                add_bar_chart(slide, data_frame, reverse_colors=reverse_colors, title_text=question,
                              sub_title_text=sub_title_text)
            elif chart_type == CHART_TYPE_COLUMN:
                add_column_chart(slide, data_frame, title_text=question, sub_title_text=sub_title_text)
            elif chart_type == CHART_TYPE_LINE:
                add_line_chart(slide, data_frame, title_text=question, sub_title_text=sub_title_text)
            elif chart_type == CHART_TYPE_PIE:

                # Rebuild data for doughnut chart
                data_frame = pd.DataFrame(columns=['First', 'Second'])

                if series_by == SERIES_BY_CATEGORIES:
                    """ Typical example for 'by categories'
                    
                    Options: Yes, No, Don’t know 
                    Section categories: Male, Female
                    
                                          First Second
                    Male                  0.804       
                    Female                0.468       
                    Male -> Yes                  0.268
                    Male -> No                   0.268
                    Male -> Don’t know           0.268
                    Female -> Yes                0.234
                    Female -> No                     0
                    Female -> Don’t know         0.234
                    """
                    for c in section['categories']:
                        data_frame.loc[c['name']] = [sum(filter(None, c['data'])), None]

                    for cat in section['categories']:
                        i = 0
                        for o in options:
                            loc_ = '{}{}{}'.format(cat['name'], PPT_CUSTOM_LABEL_SEPERATOR, o)
                            data_frame.loc[loc_] = [None, cat['data'][i]]
                            i += 1
                if series_by == SERIES_BY_OPTIONS:
                    """ Typical example for 'by options'
                    
                        Options: Yes, No, Don’t know 
                        Section categories: Male, Female
                    
                                              First Second
                        Yes                   0.502       
                        No                    0.268       
                        Don’t know            0.502       
                        Yes -> Male                  0.268
                        Yes -> Female                0.234
                        No -> Male                   0.268
                        No -> Female                     0
                        Don’t know -> Male           0.268
                        Don’t know -> Female         0.234
                    """

                    data_by_options = table['data_by_options'][cross_break]

                    i = 0
                    for o in options:
                        data_frame.loc[o] = [sum(filter(None, data_by_options[i]['data'])), None]
                        i += 1

                    i = 0
                    for o in options:
                        c = 0
                        for cat in section['categories']:
                            loc_ = '{}{}{}'.format(o, PPT_CUSTOM_LABEL_SEPERATOR, cat['name'])
                            data_frame.loc[loc_] = [None, data_by_options[i]['data'][c]]
                            c += 1
                        i += 1

                # fix table
                data_frame = (data_frame / 100).replace(np.NaN, '')

                add_doughnut_chart(slide, data_frame, title_text=question, sub_title_text=sub_title_text)

        # Strangely python-pptx's save method does not works with stream(BytesIO). Temporary file method implemented.

        save_file_name = '{}.pptx'.format(random_alphanum(20))
        save_path = os.path.join(tempfile.gettempdir(), save_file_name)

        # Save to temp file
        prs.save(save_path)

        # Read from temp file
        with open(save_path, 'rb') as f:
            output = BytesIO(f.read())

        # Remove temp file
        os.remove(save_path)

        # Return data
        return send_file(output, attachment_filename=save_file_name, as_attachment=True)


builder_bp.add_url_rule('/builder/create', view_func=Create.as_view('create'))
builder_bp.add_url_rule('/builder', view_func=Index.as_view('index'))
builder_bp.add_url_rule('/builder/<int:table_index>', view_func=TableAtIndex.as_view('table'))
builder_bp.add_url_rule('/builder/export', view_func=Export.as_view('export'))
