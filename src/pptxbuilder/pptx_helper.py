import re

import numpy as np
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import (
    XL_CHART_TYPE,
    XL_LABEL_POSITION,
    XL_LEGEND_POSITION,
    XL_TICK_MARK,
    XL_TICK_LABEL_POSITION,
    XL_MARKER_STYLE
)
from pptx.enum.text import (
    PP_ALIGN,
    MSO_ANCHOR
)
from pptx.util import (
    Pt,
    Cm
)

PPT_CHART_LEFT = Cm(0.0)
PPT_CHART_TOP = Cm(1.0)
PPT_CHART_WIDTH = Cm(25.4)
PPT_CHART_HEIGHT = height = Cm(16.3)

PPT_DEFAULT_FONT = 'Lucida Grande'

PPT_TITLE_COLOR = (51, 51, 51)
PPT_TITLE_FONT_SIZE = Pt(12)

PPT_SUB_TITLE_LEFT = Cm(0.0)
PPT_SUB_TITLE_TOP = Cm(17.6)
PPT_SUB_TITLE_WIDTH = Cm(25.4)
PPT_SUB_TITLE_HEIGHT = Cm(1.4)

PPT_SUB_TITLE_COLOR = (102, 102, 102)
PPT_SUB_TITLE_FONT_SIZE = Pt(9)

PPT_LEGEND_FONT_SIZE = Pt(8)
PPT_LEGEND_COLOR = (51, 51, 51)

PPT_CATEGORY_AXIS_LABEL_FONT_SIZE = Pt(8)
PPT_CATEGORY_AXIS_LABEL_COLOR = (102, 102, 102)

PPT_VALUE_AXIS_LABEL_FONT_SIZE = Pt(8)
PPT_VALUE_AXIS_LABEL_COLOR = (102, 102, 102)

PPT_DATA_LABEL_FONT_SIZE = Pt(8)
PPT_DATA_LABEL_COLOR = (0, 0, 0)

PPT_CUSTOM_LABEL_SEPERATOR = ' -> '


def color_palette(numofseries=0, reverse=False):
    palette = [(124, 181, 236), (67, 67, 72), (144, 237, 125), (247, 163, 92), (128, 133, 233), (241, 92, 128),
               (228, 211, 84), (43, 144, 143), (244, 91, 91), (145, 232, 225)] * 10

    if numofseries > 0:
        palette = palette[0:numofseries][::-1]

    if reverse:
        return palette[::-1]

    return palette


def strip_html_tags(text):
    """ Strip HTML tags from any string and transform special entities
    """

    rules = [
        {r'<[^<]+?>': u''},  # remove remaining tags
        {r'^\s+': u''},  # remove spaces at the beginning
        {r'\,([a-zA-Z])': r', \1'},  # add space after a comma
        {r'\s+': u' '}  # replace consecutive spaces
    ]
    for rule in rules:
        for (k, v) in rule.items():
            regex = re.compile(k)
            text = regex.sub(v, text)

    # replace special strings
    special = {
        '&nbsp;': ' ',
        '&amp;': '&',
        '&quot;': '"',
        '&lt;': '<',
        '&gt;': '>',
        '**': '',
        "’": "'"

    }
    for (k, v) in special.items():
        text = text.replace(k, v)

    return text


def clean_axes_labels(df):
    """Cleans dataframe labels. Strips html code, double white spaces and so on.
    """

    # standardise all index/column elements as unicode
    df_index_labels = df.index.map(str)
    df_col_labels = df.columns.map(str)

    col_labels = []
    index_labels = []

    for ctext in df_col_labels:
        ctext = strip_html_tags(ctext)
        col_labels.append(ctext)

    for indtext in df_index_labels:
        indtext = strip_html_tags(indtext)
        index_labels.append(indtext)

    df.columns = col_labels
    df.index = index_labels

    return df


def color_lighten(color, percent):
    color = np.array(color)
    white = np.array([255, 255, 255])
    vector = white - color
    return color + vector * percent


def chart_title(chart, title_text):
    text_frame = chart.chart_title.text_frame

    # First paragraph
    title_p = text_frame.paragraphs[0]

    # Paragraph settings
    title_p.font.color.rgb = RGBColor(*PPT_TITLE_COLOR)

    title_p.font.name = PPT_DEFAULT_FONT
    title_p.font.size = PPT_TITLE_FONT_SIZE
    title_p.font.bold = True

    title_p.text = title_text


def chart_sub_title(slide, text):
    # create textbox
    textbox = slide.shapes.add_textbox(PPT_SUB_TITLE_LEFT, PPT_SUB_TITLE_TOP, PPT_SUB_TITLE_WIDTH, PPT_SUB_TITLE_HEIGHT)

    text_frame = textbox.text_frame

    # vertical alignment
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # First paragraph
    title_p = text_frame.paragraphs[0]

    # Paragraph settings
    title_p.font.color.rgb = RGBColor(*PPT_SUB_TITLE_COLOR)
    title_p.alignment = PP_ALIGN.CENTER
    title_p.font.name = PPT_DEFAULT_FONT
    title_p.font.size = PPT_SUB_TITLE_FONT_SIZE
    title_p.font.italic = True

    title_p.text = text


def add_doughnut_chart(slide, data_frame, title_text=None, sub_title_text=None):
    excel_num_format = '0.00%'

    chart_style = 2

    # Clean data frame
    data_frame = clean_axes_labels(data_frame)

    # Adding chart data
    chart_data = ChartData()
    chart_data.categories = data_frame.index

    for i, col in enumerate(data_frame.columns):
        chart_data.add_series(col, (data_frame.ix[:, i].values), excel_num_format)

    # create chart
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT_EXPLODED, PPT_CHART_LEFT, PPT_CHART_TOP, PPT_CHART_WIDTH, PPT_CHART_HEIGHT, chart_data
    ).chart

    chart.chart_style = chart_style
    chart.has_legend = False

    if title_text is not None:
        chart_title(chart, title_text)

    if sub_title_text is not None:
        chart_sub_title(slide, sub_title_text)

    # Create main series basis color map
    # They will be used below whilst plot settings for painting points
    main_colors = color_palette()
    color_map = {}
    i = 0
    for c in chart_data._categories:
        if PPT_CUSTOM_LABEL_SEPERATOR not in c.label:
            color_map[c.label] = main_colors[i]
            i += 1

    # Python-pptx has very poor support for doughnut charts.
    # It cannot activate with standart way(chart.plots[0].has_data_labels = True)
    # I had to do iterate over chart points and change text contents of them.
    s = 0
    for serie in chart.series:
        i = 0
        for point in serie.points:
            label_ = chart_data._categories[i].label
            value_ = chart_data._series[s].values[point._idx]

            # Flag to decide whether the value should shown
            f = True

            if len(str(value_)) == 0:
                f = False

            if f:
                try:
                    value_ = float(value_)
                except ValueError:
                    f = False

            if f:
                if value_ <= 0.0:
                    f = False

            if f:
                point_text = None

                if serie.name == 'First':
                    point_text = label_
                    color = color_map[label_]

                elif serie.name == 'Second':
                    point_text = "{}\n{}".format(label_.split(PPT_CUSTOM_LABEL_SEPERATOR)[1],
                                                 "{0:.00f}%".format(value_ * 100))

                    color_key = label_.split(PPT_CUSTOM_LABEL_SEPERATOR)[0]
                    color = tuple(map(int, tuple(color_lighten(color_map[color_key], 0.2))))
                    color_map[color_key] = color

                point.data_label.text_frame.text = point_text

                # Trick to format point text
                for paragraph in point.data_label.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if serie.name == 'First':
                            run.font.size = Pt(9)
                            run.font.bold = True
                            run.font.color.rgb = RGBColor(*(247, 247, 247))
                        else:
                            run.font.size = Pt(8)

                fill_color = color
                fill = point.format.fill
                fill.solid()
                color_code = fill_color
                fill.fore_color.rgb = RGBColor(*color_code)


            else:
                point.data_label.text_frame.paragraphs[0].text = ''

            i += 1
        s += 1


def add_bar_chart(slide, data_frame, title_text=None, sub_title_text=None):
    excel_num_format = '0.00%'

    chart_style = 2

    # Clean data frame
    data_frame = clean_axes_labels(data_frame)

    # Orientation of chart type requires that we reverse the row and column order.
    data_frame = data_frame[::-1]
    data_frame = data_frame[data_frame.columns[::-1]]

    # Add chart data
    chart_data = ChartData()
    chart_data.categories = data_frame.index

    for i, col in enumerate(data_frame.columns):
        chart_data.add_series(col, (data_frame.ix[:, i].values), excel_num_format)

    # Create chart
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, PPT_CHART_LEFT, PPT_CHART_TOP, PPT_CHART_WIDTH, PPT_CHART_HEIGHT, chart_data
    ).chart

    # chart style
    chart.chart_style = chart_style

    if title_text is not None:
        chart_title(chart, title_text)

    if sub_title_text is not None:
        chart_sub_title(slide, sub_title_text)

    # Legend settings
    # Show legend if data column has more than 1 column
    if len(data_frame.columns) > 1:
        chart.has_legend = True
        legend = chart.legend

        legend.include_in_layout = False
        legend.position = XL_LEGEND_POSITION.BOTTOM
        legend.font.name = PPT_DEFAULT_FONT
        legend.font.size = PPT_LEGEND_FONT_SIZE
        legend.font.color.rgb = RGBColor(*PPT_LEGEND_COLOR)

    # Category axis (vertical) settings
    category_axis = chart.category_axis

    category_axis.visible = True
    category_axis_tick_labels = category_axis.tick_labels
    category_axis_tick_labels.offset = 400
    category_axis_tick_labels.font.name = PPT_DEFAULT_FONT
    category_axis_tick_labels.font.size = PPT_CATEGORY_AXIS_LABEL_FONT_SIZE
    category_axis_tick_labels.font.color.rgb = RGBColor(*PPT_CATEGORY_AXIS_LABEL_COLOR)

    # Value axis (horizontal) settings
    value_axis = chart.value_axis

    value_axis.visible = True
    value_axis.has_major_gridlines = True

    value_axis.major_tick_mark = XL_TICK_MARK.OUTSIDE
    value_axis.minor_tick_mark = XL_TICK_MARK.NONE
    value_axis.tick_label_position = XL_TICK_LABEL_POSITION.LOW

    value_axis_tick_labels = value_axis.tick_labels
    value_axis_tick_labels.font.size = PPT_VALUE_AXIS_LABEL_FONT_SIZE
    value_axis_tick_labels.font.name = PPT_DEFAULT_FONT
    value_axis_tick_labels.font.color.rgb = RGBColor(*PPT_VALUE_AXIS_LABEL_COLOR)
    value_axis_tick_labels.number_format = '0%'

    # Plot area settings
    plot = chart.plots[0]

    plot.overlap = -20

    # Data label settings
    plot.has_data_labels = True
    data_labels = plot.data_labels

    data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    data_labels.font.size = PPT_DATA_LABEL_FONT_SIZE
    data_labels.font.name = PPT_DEFAULT_FONT
    data_labels.font.bold = True
    data_labels.font.color.rgb = RGBColor(*PPT_DATA_LABEL_COLOR)
    data_labels.number_format = '0%'

    # Apply theme colors to series
    series_colors_list = []
    if len(data_frame.columns) >= 1:
        series_colors_list = color_palette(len(data_frame.columns))

    for i, serie in enumerate(data_frame.columns):
        serie = plot.series[i]
        serie.invert_if_negative = False

        if len(data_frame.columns) >= 1:
            fill = serie.format.fill
            fill.solid()
            color_code = series_colors_list[i]
            fill.fore_color.rgb = RGBColor(*color_code)


def add_column_chart(slide, data_frame, title_text=None, sub_title_text=None):
    # Clear data frame
    data_frame = clean_axes_labels(data_frame)

    # Create chart data
    chart_data = ChartData()
    chart_data.categories = data_frame.index

    for col in data_frame.columns:
        chart_data.add_series(col, (data_frame[col].values), '0.00%')

    # Create chart
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, PPT_CHART_LEFT, PPT_CHART_TOP, PPT_CHART_WIDTH, PPT_CHART_HEIGHT, chart_data
    ).chart

    chart.chart_style = 2

    # Titles
    if title_text is not None:
        chart_title(chart, title_text)

    if sub_title_text is not None:
        chart_sub_title(slide, sub_title_text)

    # Legend settings
    # Show legend if data column has more than 1 column
    if len(data_frame.columns) > 1:
        chart.has_legend = True
        legend = chart.legend

        legend.include_in_layout = False
        legend.position = XL_LEGEND_POSITION.BOTTOM
        legend.font.name = PPT_DEFAULT_FONT
        legend.font.size = PPT_LEGEND_FONT_SIZE
        legend.font.color.rgb = RGBColor(*PPT_LEGEND_COLOR)

    # Category axis (vertical) settings
    category_axis = chart.category_axis

    category_axis.visible = True
    category_axis_tick_labels = category_axis.tick_labels
    category_axis_tick_labels.offset = 400
    category_axis_tick_labels.font.name = PPT_DEFAULT_FONT
    category_axis_tick_labels.font.size = PPT_CATEGORY_AXIS_LABEL_FONT_SIZE
    category_axis_tick_labels.font.color.rgb = RGBColor(*PPT_CATEGORY_AXIS_LABEL_COLOR)

    # Value axis (horizontal) settings
    value_axis = chart.value_axis

    value_axis.visible = True
    value_axis.has_major_gridlines = True

    value_axis.major_tick_mark = XL_TICK_MARK.OUTSIDE
    value_axis.minor_tick_mark = XL_TICK_MARK.NONE
    value_axis.tick_label_position = XL_TICK_LABEL_POSITION.LOW

    value_axis_tick_labels = value_axis.tick_labels
    value_axis_tick_labels.font.size = PPT_VALUE_AXIS_LABEL_FONT_SIZE
    value_axis_tick_labels.font.name = PPT_DEFAULT_FONT
    value_axis_tick_labels.font.color.rgb = RGBColor(*PPT_VALUE_AXIS_LABEL_COLOR)
    value_axis_tick_labels.number_format = '0%'

    # Plot area settings
    plot = chart.plots[0]
    plot.overlap = -20

    # Data label settings
    plot.has_data_labels = True
    data_labels = plot.data_labels

    data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    data_labels.font.size = PPT_DATA_LABEL_FONT_SIZE
    data_labels.font.name = PPT_DEFAULT_FONT
    data_labels.font.bold = True
    data_labels.font.color.rgb = RGBColor(*PPT_DATA_LABEL_COLOR)
    data_labels.number_format = '0%'

    # Apply theme colors to series
    series_colors_list = []
    if len(data_frame.columns) >= 1:
        series_colors_list = color_palette(len(data_frame.columns), True)

    for i, serie in enumerate(data_frame.columns):
        serie = plot.series[i]
        serie.invert_if_negative = False

        if len(data_frame.columns) >= 1:
            fill = serie.format.fill
            fill.solid()
            color_code = series_colors_list[i]
            fill.fore_color.rgb = RGBColor(*color_code)


def add_line_chart(slide, data_frame, title_text=None, sub_title_text=None):
    # Clear data frame
    data_frame = clean_axes_labels(data_frame)

    # Create chart data
    chart_data = ChartData()
    chart_data.categories = data_frame.index

    for i, col in enumerate(data_frame.columns):
        chart_data.add_series(col, (data_frame.ix[:, i].values), '0.00%')

    # Create chart
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE, PPT_CHART_LEFT, PPT_CHART_TOP, PPT_CHART_WIDTH, PPT_CHART_HEIGHT, chart_data
    ).chart

    chart.chart_style = 2

    # Titles
    if title_text is not None:
        chart_title(chart, title_text)

    if sub_title_text is not None:
        chart_sub_title(slide, sub_title_text)

    # Legend settings
    chart.has_legend = True
    legend = chart.legend

    legend.include_in_layout = False
    legend.position = XL_LEGEND_POSITION.BOTTOM
    legend.font.name = PPT_DEFAULT_FONT
    legend.font.size = PPT_LEGEND_FONT_SIZE
    legend.font.color.rgb = RGBColor(*PPT_LEGEND_COLOR)

    # Category axis (horizontal) settings
    category_axis = chart.category_axis

    category_axis.visible = True
    category_axis_tick_labels = category_axis.tick_labels
    category_axis_tick_labels.offset = 400
    category_axis_tick_labels.font.name = PPT_DEFAULT_FONT
    category_axis_tick_labels.font.size = PPT_CATEGORY_AXIS_LABEL_FONT_SIZE
    category_axis_tick_labels.font.color.rgb = RGBColor(*PPT_CATEGORY_AXIS_LABEL_COLOR)

    # Value axis (vertical) settings
    value_axis = chart.value_axis

    value_axis.visible = True
    value_axis.has_major_gridlines = True
    value_axis.major_tick_mark = XL_TICK_MARK.OUTSIDE
    value_axis.minor_tick_mark = XL_TICK_MARK.NONE
    value_axis.tick_label_position = XL_TICK_LABEL_POSITION.LOW

    value_axis_tick_labels = value_axis.tick_labels
    value_axis_tick_labels.font.size = PPT_VALUE_AXIS_LABEL_FONT_SIZE
    value_axis_tick_labels.font.name = PPT_DEFAULT_FONT
    value_axis_tick_labels.font.color.rgb = RGBColor(*PPT_VALUE_AXIS_LABEL_COLOR)
    value_axis_tick_labels.number_format = '0%'

    # Plot area settings
    plot = chart.plots[0]

    # Data label settings
    plot.has_data_labels = True
    data_labels = plot.data_labels

    data_labels.position = XL_LABEL_POSITION.ABOVE
    data_labels.font.size = PPT_DATA_LABEL_FONT_SIZE
    data_labels.font.name = PPT_DEFAULT_FONT
    data_labels.font.bold = True
    data_labels.font.color.rgb = RGBColor(*PPT_DATA_LABEL_COLOR)
    data_labels.number_format = '0%'

    # Apply theme colors to series
    # TODO: Line chart does not apply colors
    series_colors_list = []
    if len(data_frame.columns) >= 1:
        series_colors_list = color_palette(len(data_frame.columns))

    for i, serie in enumerate(data_frame.columns):
        serie = plot.series[i]

        serie.invert_if_negative = False

        if len(data_frame.columns) >= 1:
            fill = serie.format.fill
            fill.solid()
            color_code = series_colors_list[i]
            print(color_code)
            fill.fore_color.rgb = RGBColor(*color_code)

        # Set marker style
        serie.marker.style = XL_MARKER_STYLE.AUTOMATIC
